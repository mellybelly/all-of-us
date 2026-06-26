#!/usr/bin/env python3
"""Build the AoU Roadmap Refresh draft slide deck (16:9, ~19 slides)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

HERE = os.path.dirname(os.path.abspath(__file__))
FIG  = os.path.join(HERE, "figures")
OUT  = os.path.join(HERE, "AoU_Roadmap_Refresh_Draft_Deck.pptx")

NAVY   = RGBColor(0x1A, 0x3A, 0x6B)
ACCENT = RGBColor(0x21, 0x6B, 0xB3)
DARK   = RGBColor(0x33, 0x33, 0x33)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT  = RGBColor(0xEA, 0xF0, 0xFA)
RULE   = RGBColor(0xC8, 0xC8, 0xC8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMBER  = RGBColor(0xB3, 0x59, 0x00)
SOFT_AMBER = RGBColor(0xFF, 0xF7, 0xE8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

HEAD_FONT = "Helvetica"
BODY_FONT = "Georgia"

# -- helpers ----------------------------------------------------------------

def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK,
                font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    # split paragraphs on \n
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=15, color=DARK,
                font=BODY_FONT, bullet="•  ", line_spacing=1.20,
                first_bold_prefix=False, paragraph_space=2):
    """items: list of strings, or list of (bold_prefix, rest) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(paragraph_space)
        # bullet
        rb = p.add_run()
        rb.text = bullet
        rb.font.name = font
        rb.font.size = Pt(size)
        rb.font.color.rgb = ACCENT
        rb.font.bold = True
        if isinstance(item, tuple):
            bold_part, rest = item
            if bold_part:
                rr = p.add_run()
                rr.text = bold_part
                rr.font.name = font
                rr.font.size = Pt(size)
                rr.font.bold = True
                rr.font.color.rgb = NAVY
            rr2 = p.add_run()
            rr2.text = rest
            rr2.font.name = font
            rr2.font.size = Pt(size)
            rr2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_header(slide, title, subtitle=None):
    # Top navy bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08),
             fill=NAVY, line=None)
    add_textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.55),
                title, size=22, bold=True, color=NAVY, font=HEAD_FONT,
                line_spacing=1.05)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.32),
                    subtitle, size=12, color=GRAY, font=BODY_FONT, italic=True,
                    line_spacing=1.1)
    # thin divider
    add_rect(slide, Inches(0.5), Inches(1.10), Inches(12.3), Pt(0.75),
             fill=RULE, line=None)

def add_footer(slide, slide_num, total=19):
    add_textbox(slide, Inches(0.5), Inches(7.10), Inches(8), Inches(0.3),
                "All of Us Scientific Roadmap Refresh  •  WG-1 Draft  •  Smoller & Haendel  •  2026-06-02",
                size=8.5, color=GRAY, italic=True, font=BODY_FONT)
    add_textbox(slide, Inches(11.5), Inches(7.10), Inches(1.3), Inches(0.3),
                f"{slide_num} / {total}", size=8.5, color=GRAY,
                align=PP_ALIGN.RIGHT, font=BODY_FONT, italic=True)

def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    paragraphs = text.strip().split("\n\n")
    for i, ptext in enumerate(paragraphs):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        r = p.add_run()
        r.text = ptext
        r.font.size = Pt(10)
        r.font.name = "Helvetica"

# -- table helper -----------------------------------------------------------

def add_table(slide, x, y, w, h, header, rows, *,
              header_font_size=11, body_font_size=11,
              col_widths=None, first_col_bold=False,
              last_col_emph=False, last_col_color=None):
    ncols = len(header)
    nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    table = tbl_shape.table
    if col_widths:
        total_w = w
        emu_total = total_w
        # python-pptx wants Emu widths
        used = 0
        widths_emu = []
        for i, frac in enumerate(col_widths):
            if i == ncols - 1:
                widths_emu.append(emu_total - used)
            else:
                ww = int(emu_total * frac)
                widths_emu.append(ww)
                used += ww
        for i, ww in enumerate(widths_emu):
            table.columns[i].width = ww
    # header row
    for i, ht in enumerate(header):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ht
        r.font.name = HEAD_FONT
        r.font.size = Pt(header_font_size)
        r.font.bold = True
        r.font.color.rgb = WHITE
    # body rows
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (LIGHT if ri % 2 == 1 else WHITE)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = BODY_FONT
            r.font.size = Pt(body_font_size)
            if first_col_bold and ci == 0:
                r.font.bold = True
                r.font.color.rgb = NAVY
            elif last_col_emph and ci == ncols - 1:
                r.font.bold = True
                r.font.color.rgb = last_col_color or ACCENT
                p.alignment = PP_ALIGN.CENTER
            else:
                r.font.color.rgb = DARK
    return table

# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=WHITE, line=None)
# Left navy band
add_rect(s, Inches(0), Inches(0), Inches(0.55), SLIDE_H, fill=NAVY, line=None)
# Accent stripe
add_rect(s, Inches(0.55), Inches(0), Inches(0.10), SLIDE_H, fill=ACCENT, line=None)

add_textbox(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(1.6),
            "All of Us\nScientific Roadmap Refresh",
            size=44, bold=True, color=NAVY, font=HEAD_FONT, line_spacing=1.05)
add_textbox(s, Inches(1.0), Inches(3.7), Inches(11.5), Inches(0.6),
            "Recommendations to the 2026 Science Committee",
            size=22, color=ACCENT, font=HEAD_FONT, italic=True)

# Author block
add_rect(s, Inches(1.0), Inches(5.0), Inches(6.0), Pt(1.5), fill=NAVY, line=None)
add_textbox(s, Inches(1.0), Inches(5.15), Inches(11.5), Inches(0.5),
            "Working Group 1 — Scientific Roadmap Refresh",
            size=13, bold=True, color=NAVY, font=HEAD_FONT)
add_textbox(s, Inches(1.0), Inches(5.6), Inches(11.5), Inches(0.45),
            "Co-leads:  Jordan Smoller  •  Melissa Haendel",
            size=15, color=DARK, font=BODY_FONT)
add_textbox(s, Inches(1.0), Inches(6.05), Inches(11.5), Inches(0.4),
            "Draft  •  June 2, 2026  •  for WG circulation and Science Committee pre-read",
            size=11, color=GRAY, font=BODY_FONT, italic=True)

set_notes(s, """Title slide for the 30-minute WG-1 presentation to the AoU Science Committee.

Co-leads Smoller and Haendel will jointly present. Plan: ~12 minutes of evidence and findings, ~12 minutes of recommendations, ~6 minutes of Q&A and decision-asks.

Audience: Science Committee, AoU program leadership (Schully, Roden, Lin), NIH program staff observers, and the other four 2026 WG leads who will fit their recommendations into whatever framework this WG endorses.

Companion artifact: AoU_Roadmap_Refresh_Executive_Summary.pdf (one-page Sci Com pre-read).""")

# ---------------------------------------------------------------------------
# SLIDE 2 — Charge
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Our charge: refresh the Roadmap against Roadmap 2.0's six questions",
           "WG-1 owns framework changes that no other WG can make. We must answer them so WG-2/3/4/5 have a coherent home.")

add_bullets(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5), [
    ("What has changed since 2023? ",
     "— evidence from the AoU corpus and the field's scientific direction."),
    ("What are the biggest data gaps? ",
     "— where IC strategic plans demand capability AoU isn't yet producing."),
    ("How do we ensure research reproducibility? ",
     "— framework, tooling, and culture expectations for the next 5 years."),
    ("What new tools are needed? ",
     "— researcher-facing capability the Workbench must add or enable."),
    ("What exposures and outcomes should we capture, and how do we demonstrate impact? ",
     "— scope, instruments, and translational-impact measurement."),
    ("Are the 8 priority focus areas still right? ",
     "— structural review of the framework. ONLY this WG can change it."),
], size=17, line_spacing=1.30, paragraph_space=8)

add_rect(s, Inches(0.6), Inches(6.30), Inches(12.1), Inches(0.65), fill=LIGHT, line=NAVY)
add_textbox(s, Inches(0.75), Inches(6.38), Inches(11.8), Inches(0.5),
            "Strategic implication: only WG-1 can change the framework. Imaging, AI, and Communications recommendations must fit inside whatever we decide.",
            size=12.5, color=NAVY, italic=True, bold=True, font=BODY_FONT, line_spacing=1.2)

add_footer(s, 2)
set_notes(s, """Source: slide 8 of '2026 Science Committee Working Groups.pptx' (the six Roadmap 2.0 Key Questions).

The structural-review question (Q6) is the headline because focus-area structure is uniquely WG-1's responsibility — WG-2 (Translational), WG-3 (Imaging), WG-4 (AI), and WG-5 (Communications) all need to know what the framework looks like before they can finalize their own recommendations.

The other five questions don't disappear; they get folded into the evidence base.""")

# ---------------------------------------------------------------------------
# SLIDE 3 — Approach
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "We took a five-step, evidence-driven approach",
           "Every recommendation here can be traced to a primary corpus, IC plan, peer biobank, or rubric cell.")

steps = [
    ("1", "Inventory", "the AoU corpus — 1,428 substantive publications + 12,899 substantive workspace projects, theme-tagged against a locked 19-theme taxonomy."),
    ("2", "Map", "AoU themes onto 29 NIH strategic plans (NIH-Wide + 21 ICs + cross-cutting). Extracted 166 priority statements; built a corpus × IC alignment matrix."),
    ("3", "Compare", "AoU corpus against 4 peer biobanks (UK Biobank, FinnGen, MVP, China Kadoorie) using a same-tagger keyword approach — calibrates AoU's distinctive contributions and gaps."),
    ("4", "Identify", "structural framework gaps, emerging research themes, and demand-supply mismatches across the 8 focus areas."),
    ("5", "Score", "candidate recommendations against the 8-criterion Sci Com rubric (the format Smoller established in November 2024) and rank the Top-3 per WG."),
]

# Numbered step circles + text
y0 = Inches(1.45)
step_h = Inches(1.10)
for i, (num, verb, rest) in enumerate(steps):
    y = y0 + step_h * i
    # numbered circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.85), Inches(0.85))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
    circ.line.color.rgb = NAVY
    tf = circ.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.bold = True
    r.font.size = Pt(22); r.font.color.rgb = WHITE; r.font.name = HEAD_FONT
    # verb + rest
    add_textbox(s, Inches(1.7), y + Inches(0.05), Inches(11.1), Inches(0.4),
                verb, size=15, bold=True, color=ACCENT, font=HEAD_FONT,
                line_spacing=1.05)
    add_textbox(s, Inches(1.7), y + Inches(0.42), Inches(11.1), Inches(0.7),
                rest, size=12.5, color=DARK, font=BODY_FONT, line_spacing=1.20)

add_footer(s, 3)
set_notes(s, """Phase numbering inside the evidence base:
- Phase 1: corpus acquisition + theme tagging
- Phase 2: validation + audit
- Phase 3a: 19→8 crosswalk
- Phase 3b: temporal trends pre/post-2023
- Phase 3c: IC alignment + gap analysis
- Phase 4: top-3 per WG with rubric scoring

The rubric is Smoller's 8-criterion format from the 11/5/24 Pre-Processed Data subgroup deck. Worked-example scores landed mean 1.40–4.00; most are 3.0–3.8.""")

# ---------------------------------------------------------------------------
# SLIDE 4 — Evidence base by the numbers
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Evidence base: numbers that anchor every recommendation",
           "Where claims rest on small samples (peer biobanks, pre-2023 baseline), we flag the limitation explicitly.")

stats = [
    ("1,428",     "substantive AoU\npublications",       NAVY),
    ("12,899",    "substantive workspace\nprojects",     NAVY),
    ("29",        "NIH strategic plans\n(NIH-Wide + 21 ICs)", NAVY),
    ("166",       "IC strategic priorities\nextracted",  NAVY),
    ("4",         "peer-biobank corpora\n(UKB, FinnGen, MVP, CKB)", NAVY),
    ("19",        "themes in the\nlocked taxonomy",      NAVY),
    ("8",         "current AoU Roadmap\nfocus areas",    NAVY),
    ("8",         "Sci Com rubric criteria\n(Smoller, Nov 2024)", NAVY),
]

# 2 rows × 4 columns
card_w = Inches(2.95); card_h = Inches(2.10)
gap_x = Inches(0.12); gap_y = Inches(0.20)
start_x = Inches(0.55); start_y = Inches(1.50)

for i, (big, label, color) in enumerate(stats):
    row = i // 4; col = i % 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect(s, x, y, card_w, card_h, fill=LIGHT, line=NAVY)
    # accent stripe top
    add_rect(s, x, y, card_w, Inches(0.10), fill=ACCENT, line=None)
    add_textbox(s, x, y + Inches(0.25), card_w, Inches(1.0),
                big, size=42, bold=True, color=NAVY, font=HEAD_FONT,
                align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_textbox(s, x + Inches(0.1), y + Inches(1.35), card_w - Inches(0.2),
                Inches(0.7), label, size=11, color=DARK, font=BODY_FONT,
                align=PP_ALIGN.CENTER, line_spacing=1.15)

add_footer(s, 4)
set_notes(s, """Corpus tally is Phase 2c audited substantive corpus. Earlier intake counts included program-meta and methods-not-yet-applied workspaces; we excluded those.

The 4 peer biobanks were pulled via PubMed name-search with same-tagger keyword tagging — comparisons are internally consistent but absolute counts should not be treated as authoritative.

The 8-criterion rubric: (1) unmet researcher need; (2) simplifies analysis; (3) standardizes/harmonizes; (4) broad use cases; (5) Scientific Roadmap alignment; (6) tool exists vs new development; (7) no new data collection; (8) helps with data quality. Each scored 1–4 (1=low, 4=high); N/A excluded.""")

# ---------------------------------------------------------------------------
# SLIDE 5 — 19-theme taxonomy at a glance
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "The 19-theme taxonomy: where AoU science actually lives",
           "Cardiometabolic, Genomics, and Methods/AI together make up ~70% of the corpus.")

# Two-column theme list with pub-share
theme_data = [
    ("1", "Cardiometabolic disease",                            "28.4%", True),
    ("2", "Genomics methods & complex-trait genetics",          "24.9%", True),
    ("3", "Methods, infrastructure & phenotyping",              "18.4%", True),  # highlighted
    ("4", "Mental health, behavioral & substance use",          "18.8%", False),
    ("5", "Cancer",                                              "12.8%", False),
    ("6", "Social determinants & health disparities",           "32.9%", False),
    ("7", "Wearables & digital health",                          "5.4%",  False),
    ("8", "Infectious disease & immunology",                     "6.0%",  False),
    ("9", "Autoimmune & inflammatory disease",                   "7.6%",  False),
    ("10", "Rare disease & Mendelian genetics",                  "3.8%",  False),
    ("11", "Aging, frailty & ADRD",                              "5.5%",  False),
    ("12", "Dermatology (artifact — flag)",                      "13.9%", False),
    ("13", "Respiratory & sleep medicine",                       "4.2%",  False),
    ("14", "Neurology (non-ADRD)",                               "5.7%",  False),
    ("15", "Pregnancy, maternal & women's health",               "2.7%",  False),
    ("16", "Pharmacogenomics & drug response",                   "2.4%",  False),
    ("17", "ELSI, participant engagement & recruitment",         "4.9%",  False),
    ("18", "Ophthalmology, ENT & sensory medicine",              "7.6%",  False),
    ("19", "Environmental health & climate",                     "0.6%",  False),
]

col_w = Inches(6.15); col_h = Inches(0.27); start_y = Inches(1.40)
left_x = Inches(0.55); right_x = Inches(6.85)

# Header for the two columns
for cx in (left_x, right_x):
    add_textbox(s, cx, Inches(1.18), col_w, Inches(0.25),
                "# THEME                                                                      SHARE",
                size=9, bold=True, color=NAVY, font=HEAD_FONT,
                line_spacing=1.0)

for i, (num, name, share, big) in enumerate(theme_data):
    col = 0 if i < 10 else 1
    row = i if i < 10 else i - 10
    x = left_x if col == 0 else right_x
    y = start_y + col_h * row
    fill = SOFT_AMBER if big else (LIGHT if row % 2 == 0 else WHITE)
    add_rect(s, x, y, col_w, col_h, fill=fill, line=None)
    # num
    add_textbox(s, x + Inches(0.06), y, Inches(0.35), col_h,
                num, size=10, bold=True, color=NAVY, font=HEAD_FONT,
                anchor=MSO_ANCHOR.MIDDLE)
    # name
    color_n = AMBER if big else DARK
    add_textbox(s, x + Inches(0.5), y, Inches(4.6), col_h,
                name, size=10.5, color=color_n, font=BODY_FONT,
                bold=big, anchor=MSO_ANCHOR.MIDDLE)
    # share
    add_textbox(s, x + col_w - Inches(0.9), y, Inches(0.85), col_h,
                share, size=10.5, color=(ACCENT if big else GRAY),
                bold=big, font=HEAD_FONT, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)

# Legend / annotation
add_textbox(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.45),
            "Highlighted: Theme 3 (Methods/AI) is the single largest corpus orphan — currently with NO home in the 8-area framework.",
            size=11, color=AMBER, italic=True, font=BODY_FONT, bold=True)
add_textbox(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.30),
            "Shares are share-of-corpus-touched (pubs can multi-tag); rows therefore sum >100%.",
            size=9, color=GRAY, italic=True, font=BODY_FONT)

add_footer(s, 5)
set_notes(s, """Source: aou-research-landscape/theme_summary.csv

Theme 3 (Methods/AI) — 262 pubs / 2,388 projects, 18.4% pub share / 18.5% project share — is the single largest orphan in the 19→8 crosswalk. This is the empirical case for the 9th focus area.

Theme 12 (Dermatology) at 13.9% share is anomalous and driven by 1–2 prolific groups per Phase 1 README; trend is declining (−4.4). Flagged on slide 18 (Caveats).

Theme 6 (SDOH) at 33.3% is the largest single theme — defines AoU as a cohort.""")

# ---------------------------------------------------------------------------
# SLIDE 6 — Q6: 8 focus areas still right
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "The 8 focus areas mostly hold — but three structural gaps need fixing",
           "Q6 answer: yes at the top level, with three named exceptions only WG-1 can address.")

# Three problem cards
gap_y = Inches(1.4)
card_w = Inches(4.07); card_h = Inches(4.3); gap_x = Inches(0.12)
sx = Inches(0.55)

cards = [
    ("ORPHAN", "Methods, AI & Infrastructure",
     "18% of substantive pubs (262) and 19% of projects (2,388) have NO clean home in the 8-area framework.",
     "Includes the AI/foundation-model surge (HA4: 33 pubs, 30 post-2023), workbench tooling, EHR phenotyping, federated learning.",
     "Sub-signal: AoU peer-index 3.0× — distinctive AoU strength currently invisible in the framework."),
    ("ORPHAN", "Wearables / remote monitoring",
     "75 pubs (5.4% share) — slide-7 priority that has nowhere to live in the 8 areas.",
     "Trend is flat (+0.2) — wearables have already MATURED, not emerging. Important reframe.",
     "Stretch-fit into FA5 (Behavioral) misses sleep, HRV, circadian, multimodal EHR+wearable AI work."),
    ("INVISIBLE", "Multi-omics inside \"Genetics & Biology\"",
     "Multi-omics is a slide-7 2023+ priority that is buried inside FA8.",
     "HA1: 24 pubs (23 post-2023), mostly external proteomics linked to AoU genotypes.",
     "Largest peer-comparative growth differential: UKB 12.2% post-2023 vs AoU 1.8%. 2026 omics releases need a NAMED home."),
]

for i, (badge, title, l1, l2, l3) in enumerate(cards):
    x = sx + i * (card_w + gap_x)
    add_rect(s, x, gap_y, card_w, card_h, fill=WHITE, line=NAVY, line_w=1.2)
    add_rect(s, x, gap_y, card_w, Inches(0.45), fill=NAVY, line=None)
    add_textbox(s, x + Inches(0.15), gap_y + Inches(0.05), card_w, Inches(0.35),
                badge, size=11, bold=True, color=WHITE, font=HEAD_FONT)
    add_textbox(s, x + Inches(0.15), gap_y + Inches(0.55), card_w - Inches(0.3),
                Inches(0.9), title, size=17, bold=True, color=NAVY, font=HEAD_FONT,
                line_spacing=1.10)
    add_bullets(s, x + Inches(0.1), gap_y + Inches(1.50),
                card_w - Inches(0.2), card_h - Inches(1.5),
                [l1, l2, l3], size=11.5, line_spacing=1.25, paragraph_space=5)

# Recommendation strip
add_rect(s, Inches(0.55), Inches(5.9), Inches(12.25), Inches(0.95), fill=SOFT_AMBER, line=AMBER, line_w=1.2)
add_textbox(s, Inches(0.75), Inches(5.95), Inches(11.85), Inches(0.4),
            "WG-1 recommendation (Rec 1): add a 9th focus area — \"Data, AI & Methods\"",
            size=14, bold=True, color=AMBER, font=HEAD_FONT)
add_textbox(s, Inches(0.75), Inches(6.30), Inches(11.85), Inches(0.55),
            "Absorbs theme 3 (the 18% orphan), wearables, AI/foundation models, and imaging-as-data once it flows. Without it, WG-3 (Imaging) and WG-4 (AI) recommendations have no coherent home in the Roadmap.",
            size=11.5, color=DARK, font=BODY_FONT, italic=True, line_spacing=1.25)

add_footer(s, 6)
set_notes(s, """Source: aou-roadmap-refresh/crosswalk_19_to_8.md §4

The crosswalk's Q6 answer in one paragraph: the 8 areas hold for disease anchoring (FA1 74%, FA8 54%, FA7 50%) but three structural problems exist — (a) theme 3 methods/AI orphan, (b) wearables orphan despite being a slide-7 priority, (c) multi-omics + imaging + AI buried in FA8 / theme 3.

The 9th focus area is the single most important recommendation in the entire WG package — it's the one change ONLY this WG can make, and three downstream WGs (Imaging, AI, Communications) depend on it.""")

# ---------------------------------------------------------------------------
# SLIDE 7 — Q1: What's changed since 2023 — already realized
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Four slide-7 priorities are already realized in the corpus",
           "Q1 (what's changed since 2023): the AoU researcher base is delivering on these without further investment.")

realized = [
    ("AI / LLMs / foundation models",
     "33 pubs (HA4); 30 post-2023",
     "PMID:39794311 foundation-model phenome-wide onset prediction; PMID:40417537 DL time-to-event depression/asthma; PMID:40750633 PAL-AI oocyte maturation.",
     "Theme 3 peer-index 3.0×"),
    ("Breakthrough therapeutics",
     "20 pubs (HA8); 19 post-2023",
     "GLP-1 hepatic decompensation (PMID:41847743); third-line CRC cost-effectiveness (PMID:41807852); obesity meds vs bariatric trends (PMID:41741949).",
     "Clearest 'translation in progress' signal in the corpus"),
    ("Wearables / remote monitoring",
     "75 pubs (HA2); 67 post-2023; trend flat",
     "Mature, not emerging. Sleep-stage, HRV, circadian, multimodal EHR+wearable AI work all happening — but invisible in the focus-area framework.",
     "Reframe: recognition gap, not growth gap"),
    ("G×E×SDOH integrated holistic risk",
     "11 pubs (HA6); 0 pre-2023",
     "Genuinely novel category. PMID:41518215 PRS + neighborhood disorder + sleep; PMID:42109163 CoQ10 PRS + statin myopathy.",
     "AoU positioned to LEAD the field here"),
]

for i, (title, count, examples, takeaway) in enumerate(realized):
    row = i // 2; col = i % 2
    x = Inches(0.55) + col * Inches(6.20)
    y = Inches(1.40) + row * Inches(2.78)
    add_rect(s, x, y, Inches(6.05), Inches(2.65), fill=LIGHT, line=NAVY)
    add_rect(s, x, y, Inches(6.05), Inches(0.45), fill=ACCENT, line=None)
    add_textbox(s, x + Inches(0.15), y + Inches(0.05), Inches(5.85), Inches(0.4),
                title, size=14, bold=True, color=WHITE, font=HEAD_FONT)
    add_textbox(s, x + Inches(0.18), y + Inches(0.55), Inches(5.75), Inches(0.35),
                count, size=11.5, bold=True, color=NAVY, font=HEAD_FONT, italic=True)
    add_textbox(s, x + Inches(0.18), y + Inches(0.95), Inches(5.75), Inches(1.15),
                examples, size=10.5, color=DARK, font=BODY_FONT, line_spacing=1.25)
    add_textbox(s, x + Inches(0.18), y + Inches(2.10), Inches(5.75), Inches(0.5),
                "→ " + takeaway, size=11, color=ACCENT, italic=True, bold=True,
                font=BODY_FONT, line_spacing=1.20)

add_footer(s, 7)
set_notes(s, """Source: aou-roadmap-refresh/temporal_trends.md §4–§5; phase4 cross_cutting_synthesis.md §3.

These four areas are RESEARCHER-DRIVEN and need recognition, not investment. The framework change (9th focus area) is what makes them visible.

Wearables specifically: the trend score is +0.2 — flat in relative terms. Wearables have already matured to a permanent place in the work but the framework hasn't caught up. This is a 'recognition gap, not growth gap' — important to surface explicitly so Sci Com doesn't try to invest in something that's already steady-state.""")

# ---------------------------------------------------------------------------
# SLIDE 8 — Q1: aspirational, pending upstream investment
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Four slide-7 priorities are aspirational — pending upstream investment",
           "Q1 (continued): these need data-supply or sample-collection decisions, not researcher engagement.")

aspirational = [
    ("Multi-omics beyond genomics",
     "24 pubs (HA1); 23 post-2023",
     "Mostly external proteomics/metabolomics linked to AoU genotypes (PMID:41712304 COPD proteomics+PRS; PMID:40050615 fibroid GWAS+eQTL).",
     "UKB 12.2% post-2023 / AoU 1.8% — largest growth differential. 2026 release closes mechanically."),
    ("Imaging / radiomics",
     "8 pubs total (HA3) — below thin-coverage threshold",
     "All 8 use EXTERNAL imaging linked to AoU phenotypes (UKB MRI, EHR mammography). No AoU-originated imaging.",
     "AoU corpus shows almost no imaging delivery yet. WG-3 starts from infrastructure."),
    ("Nutrition / microbiome",
     "11 pubs (HA7); structurally limited",
     "AoU dietary surveys thin vs UKB 24-hour recall; no stool biospecimens collected.",
     "UKB 11.3% / FinnGen 14.6% post-2023 / AoU 0.8%. Requires sample-collection decision."),
    ("Environmental health & climate",
     "13 pubs (HA5); 8 substantive (theme 19)",
     "Largest single demand-supply gap in the IC alignment matrix. NIEHS has 6 plan priorities tagged here.",
     "Requires PM2.5 / EJI / heat / water linkage investment — engineering, not data acquisition."),
]

for i, (title, count, examples, takeaway) in enumerate(aspirational):
    row = i // 2; col = i % 2
    x = Inches(0.55) + col * Inches(6.20)
    y = Inches(1.40) + row * Inches(2.78)
    add_rect(s, x, y, Inches(6.05), Inches(2.65), fill=SOFT_AMBER, line=AMBER)
    add_rect(s, x, y, Inches(6.05), Inches(0.45), fill=AMBER, line=None)
    add_textbox(s, x + Inches(0.15), y + Inches(0.05), Inches(5.85), Inches(0.4),
                title, size=14, bold=True, color=WHITE, font=HEAD_FONT)
    add_textbox(s, x + Inches(0.18), y + Inches(0.55), Inches(5.75), Inches(0.35),
                count, size=11.5, bold=True, color=AMBER, font=HEAD_FONT, italic=True)
    add_textbox(s, x + Inches(0.18), y + Inches(0.95), Inches(5.75), Inches(1.15),
                examples, size=10.5, color=DARK, font=BODY_FONT, line_spacing=1.25)
    add_textbox(s, x + Inches(0.18), y + Inches(2.10), Inches(5.75), Inches(0.5),
                "→ " + takeaway, size=11, color=AMBER, italic=True, bold=True,
                font=BODY_FONT, line_spacing=1.20)

add_footer(s, 8)
set_notes(s, """Source: aou-roadmap-refresh/temporal_trends.md §4–§5; peer-biobanks/emerging_themes.md §1–§3.

The clean split (4 realized + 4 aspirational) reframes the Q1 conversation. Instead of 'are we delivering on the 2023 priorities, yes or no?', the answer is: the researcher-driven priorities ARE delivering; the data-supply-dependent priorities are not, by construction, and need infrastructure decisions before they can.

Environmental is the canonical demand-side gap: NIEHS has 6 plan priorities; AoU has 7 substantive pubs. This is the strongest evidence for the FA6 environmental-linkage commitment in Rec 3.""")

# ---------------------------------------------------------------------------
# SLIDE 9 — Strongest IC alignments
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Three IC alignments stand out as AoU's strongest current cases",
           "Each is anchored on plain-text plan language we can quote in funding conversations.")

ic_rows = [
    ("NHGRI × Genomics methods",
     "Largest single matrix cell. AoU delivers Bold Prediction #9 (equitable benefit from genomics for ancestrally diverse populations).",
     "\"…the equitable benefit of genomics for all populations.\" — NHGRI 2020 Strategic Vision, BP #9"),
    ("NIMHD × SDOH / disparities",
     "Goal 7 literally describes the AoU mandate. AoU 32.9% SDOH share; 5.1× peer median.",
     "\"…analysis of populations experiencing health disparities in big data sets, clinical research, and future big science initiatives.\" — NIMHD G7"),
    ("NLM × Methods / infrastructure",
     "The only IC plan that explicitly names All of Us. Theme 3: 262 pubs, 2,388 projects.",
     "\"NLM will engage with the entire NIH to determine how to accommodate data generated in … All of Us.\" — NLM 2017–2027 Strategic Plan"),
]

t = add_table(s, Inches(0.55), Inches(1.40), Inches(12.25), Inches(4.5),
              header=["IC × Theme", "Why it's strong", "Anchor language we can cite"],
              rows=ic_rows,
              header_font_size=12, body_font_size=11.5,
              col_widths=[0.22, 0.36, 0.42],
              first_col_bold=True)

# Row height tuning
for r in t.rows:
    r.height = Inches(1.45)
t.rows[0].height = Inches(0.45)

add_rect(s, Inches(0.55), Inches(6.10), Inches(12.25), Inches(0.85),
         fill=LIGHT, line=NAVY)
add_textbox(s, Inches(0.75), Inches(6.18), Inches(11.85), Inches(0.7),
            "Headline read: lead the WG-1 narrative with the genomics-and-genomics story; NIMHD G7 is the strongest political-defense anchor; NLM is the only mission-explicit partner.",
            size=12.5, color=NAVY, italic=True, bold=True, font=BODY_FONT, line_spacing=1.25)

add_footer(s, 9)
set_notes(s, """Source: aou-ic-alignment/gap_analysis.md §A (alignment hits) + per_ic_briefs/.

NHGRI is the strongest current IC home for the genetics/multi-omics framing (Rec 2). NHGRI Bold Prediction #5 (millions of human participants) implicitly requires AoU; #9 (equitable benefit) is AoU's literal contribution.

NIMHD G7 is the strongest 'AoU embodies a pan-NIH mandate' argument — the language is so specific to AoU's design that it reads as if written about AoU.

NLM is the only IC plan that names AoU; NLM-G1-O2 (informatics R&D) and NLM-G1-O4 (sustainable computational infrastructure) directly underwrite the 9th focus area.""")

# ---------------------------------------------------------------------------
# SLIDE 10 — Biggest demand-side gaps
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Three demand-side gaps: ICs want capability AoU isn't yet producing",
           "Each requires a different kind of decision (engineering, governance, or sample collection).")

gaps = [
    ("NIEHS × Environment",
     "6 IC priorities  ↔  7 AoU substantive pubs",
     "LARGEST demand-supply gap in the entire matrix.",
     "Engineering decision",
     "PM2.5, EJI, heat, water — external public data + AoU geocoding."),
    ("NCI / NEI / NIBIB × Imaging",
     "Multiple IC priorities  ↔  8 AoU pubs total",
     "Data-supply problem (no AoU imaging release), not researcher demand problem.",
     "Governance + release decision",
     "Tier 1 retinal/ECG/spirometry; Tier 2 EHR DICOM; Tier 3 prospective."),
    ("NIDDK / NHLBI × Multi-omics",
     "NIDDK RO 1.3 + NHLBI Obj 4  ↔  24 AoU pubs",
     "Mostly external proteomics linked to AoU genotypes today.",
     "Closes mechanically — 2026 omics release",
     "Plus FA8 rename to send a forward signal (Rec 2)."),
    ("NIA × Aging / ADRD",
     "Lecanemab/donanemab era  ↔  AoU at 6.0% (near peer median)",
     "Peer growth (UKB +6.9; MVP +5.4; CKB +5.1 post-2023); AoU growing (+4.7) but trailing UKB despite structural advantage.",
     "Engagement decision",
     "Surfaces as cross-cutting flag (no single WG owns it; CC-5)."),
]

for i, (ic, ratio, problem, decision_type, action) in enumerate(gaps):
    row = i // 2; col = i % 2
    x = Inches(0.55) + col * Inches(6.20)
    y = Inches(1.40) + row * Inches(2.78)
    add_rect(s, x, y, Inches(6.05), Inches(2.65), fill=WHITE, line=NAVY, line_w=1.2)
    add_rect(s, x, y, Inches(6.05), Inches(0.45), fill=AMBER, line=None)
    add_textbox(s, x + Inches(0.15), y + Inches(0.05), Inches(5.85), Inches(0.4),
                ic, size=13.5, bold=True, color=WHITE, font=HEAD_FONT)
    add_textbox(s, x + Inches(0.18), y + Inches(0.55), Inches(5.75), Inches(0.35),
                ratio, size=11.5, bold=True, color=NAVY, font=HEAD_FONT)
    add_textbox(s, x + Inches(0.18), y + Inches(0.93), Inches(5.75), Inches(0.7),
                problem, size=10.5, color=DARK, font=BODY_FONT, line_spacing=1.22)
    add_textbox(s, x + Inches(0.18), y + Inches(1.70), Inches(5.75), Inches(0.4),
                "Decision type: " + decision_type, size=10.5,
                color=ACCENT, font=HEAD_FONT, bold=True, italic=True)
    add_textbox(s, x + Inches(0.18), y + Inches(2.05), Inches(5.75), Inches(0.5),
                "→ " + action, size=10.5, color=AMBER, italic=True, bold=True,
                font=BODY_FONT, line_spacing=1.20)

add_footer(s, 10)
set_notes(s, """Source: aou-ic-alignment/gap_analysis.md §B (demand-side gaps).

The four gap categories require different decision types — important to surface for Sci Com so the conversation doesn't conflate them:

- Environment: engineering work (external data + geocoding). LOW cost, HIGH yield.
- Imaging: governance + release decision. HIGH cost. WG-3 owns.
- Multi-omics: mechanical closure via 2026 releases. Framework signal (FA8 rename) is the WG-1 lever.
- Aging/ADRD: engagement decision, not infrastructure. Cross-cutting flag, surfaces in WG-5 patient-facing comms.""")

# ---------------------------------------------------------------------------
# SLIDE 11 — Peer-biobank comparison
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Peer-biobank comparison: AoU's gaps and AoU's distinctive strengths",
           "Same-tagger keyword search across UK Biobank (11,868 pubs), FinnGen (3,152), MVP (438), and CKB (558) — full PubMed census.")

# Two-column: under-indexed and over-indexed
add_textbox(s, Inches(0.55), Inches(1.30), Inches(6.0), Inches(0.4),
            "Where AoU is UNDER-INDEXED vs peers",
            size=14, bold=True, color=AMBER, font=HEAD_FONT)

# Use the peer_under_indexed figure
s.shapes.add_picture(os.path.join(FIG, "peer_under_indexed.png"),
                     Inches(0.45), Inches(1.65), width=Inches(6.4))

add_textbox(s, Inches(7.05), Inches(1.30), Inches(6.0), Inches(0.4),
            "Where AoU LEADS the peer median",
            size=14, bold=True, color=ACCENT, font=HEAD_FONT)
s.shapes.add_picture(os.path.join(FIG, "distinctive_strengths.png"),
                     Inches(6.93), Inches(1.65), width=Inches(6.35))

# Bottom narrative
add_rect(s, Inches(0.55), Inches(6.20), Inches(12.25), Inches(0.85),
         fill=LIGHT, line=NAVY)
add_textbox(s, Inches(0.75), Inches(6.28), Inches(11.85), Inches(0.4),
            "The narrative sentence (for the Sci Com presentation and WG-5 ambassador toolkit):",
            size=11, bold=True, color=NAVY, italic=True, font=HEAD_FONT)
add_textbox(s, Inches(0.75), Inches(6.55), Inches(11.85), Inches(0.5),
            "\"Every other cohort scaled up genomics; AoU scaled up everything around genomics that genomics actually needs to work in practice.\"",
            size=11.5, color=NAVY, italic=True, bold=True, font=BODY_FONT, line_spacing=1.2)

add_footer(s, 11)
set_notes(s, """Source: peer-biobanks/aou_vs_peers_index.csv + emerging_themes.md.

Under-indexed (left chart): multi-omics 1.7% vs UKB 10%, imaging 0.6% vs UKB 8.0%, microbiome 0.8% vs FinnGen 14%, rare disease 3.8% vs UKB 18%, aging 5.5% vs UKB 15.7%.

Over-indexed (right chart): ELSI 26.3× peer median, SDOH 5.1×, PGx 4.6×, methods 3.0×. These are AoU's DEFENSIBLE IDENTITY — distinctive strengths the Refresh must preserve.

Caveats: PubMed-name-search has known undercount (especially MVP); pre-2023 baseline is small (157 AoU pubs); use STATIC SHARE INDEX, not trend differentials. Dermatology 6.9× over-index is an artifact (1-2 prolific groups).""")

# ---------------------------------------------------------------------------
# SLIDE 12 — WG-1 (OURS) Top-3
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "WG-1 Top-3: the framework changes only this WG can make",
           "Rubric scores in parentheses. Each row includes the concrete decision we ask the Sci Com to make.")

# Table format for WG-1 with extra detail
rows = [
    ("1",
     "Add a 9th focus area: \"Data, AI & Methods\"",
     "Absorbs theme 3 (18% orphan), wearables (HA2 75 pubs), AI/foundation models (HA4 33 pubs), and imaging-as-data once it flows.",
     "3.71",
     "ENDORSE the framework addition"),
    ("2",
     "Rename FA8 → \"Genetics, Multi-omics & Biology\" + commit to proteomics roadmap",
     "Signal that 2026 omics releases (Olink / SomaScan) have a home. UKB 12.2% post-2023 / AoU 1.8% — largest growth differential.",
     "3.38",
     "ENDORSE rename + AUTHORIZE release roadmap"),
    ("3",
     "Bundle FA reframes: FA4 toward ROR implementation; FA2 → \"Maternal Health\"; FA6 environmental-exposure linkage commitment",
     "FA4 trend −9.6 (ELSI work maturing out, ROR implementation continues); FA2 has no Child content; FA6 is the largest demand-supply gap.",
     "3.00",
     "DISCUSS + vote on the three sub-items"),
]

t = add_table(s, Inches(0.55), Inches(1.35), Inches(12.25), Inches(4.7),
              header=["#", "Recommendation", "Evidence",
                      "Rubric", "Action requested"],
              rows=rows,
              header_font_size=12, body_font_size=11,
              col_widths=[0.04, 0.30, 0.39, 0.09, 0.18],
              last_col_emph=True, last_col_color=AMBER)
# Row heights
t.rows[0].height = Inches(0.42)
for ri in range(1, 4):
    t.rows[ri].height = Inches(1.42)

# Why this matters
add_rect(s, Inches(0.55), Inches(6.20), Inches(12.25), Inches(0.85),
         fill=SOFT_AMBER, line=AMBER)
add_textbox(s, Inches(0.75), Inches(6.28), Inches(11.85), Inches(0.4),
            "Why these and not more?  The Sci Com rubric format limits us to three. Each here is structural; each is uniquely WG-1.",
            size=11.5, bold=True, color=AMBER, font=HEAD_FONT)
add_textbox(s, Inches(0.75), Inches(6.55), Inches(11.85), Inches(0.5),
            "Aging/ADRD, microbiome, reproducibility framework, subcontinental ancestry — surfaced in cross-cutting synthesis as Tier-2 or other-WG.",
            size=11, color=DARK, font=BODY_FONT, italic=True, line_spacing=1.20)

add_footer(s, 12)
set_notes(s, """Source: phase4-recommendations/wg1_roadmap_refresh.md (Top-3 with full rubric tables).

Rec 1 (3.71 mean): 'Add a Data, AI & Methods focus area.' 18% corpus orphan; HA4 33 AI pubs; NLM-G1 explicit alignment. Risk: Sci Com may prefer reorganization to addition; defend with the empirical orphan count.

Rec 2 (3.38 mean): 'Rename FA8 + multi-omics roadmap.' Largest peer-comparative growth differential. The rename without the data is a hollow gesture; the data without the rename buries the work — recommendation explicitly pairs them.

Rec 3 (3.00 mean): three small structural fixes bundled. FA4 toward ROR implementation (well-evidenced; theme 17 ELSI −12.3); FA2 → 'Maternal Health' (judgment about naming accuracy; politically sensitive — frame as accuracy not deprioritization); FA6 environmental linkage (largest demand-supply gap; engineering not data acquisition).""")

# ---------------------------------------------------------------------------
# SLIDE 13 — WG-2 Translational Top-3
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "WG-2 Translational: define what counts as translational impact",
           "Lead: Sanchez & Cohn. Recommendations bridge WG-1's framework and WG-5's reporting.")

rows = [
    ("1", "Adopt a formal translational-impact metric set (TSBM-based)",
     "Replication in peers; consortium adoption (RDCRN, ClinGen, GREGoR, eMERGE); FDA/CMS citations; NCT enablement; CDC/USPSTF/NAM citations.",
     "3.50"),
    ("2", "Name 3–5 priority translational use cases",
     "ACMG actionable-gene implementation (54 pubs/433 projects); multi-ancestry PRS clinical pilots; PGx implementation (4.6× peer); cardiometabolic+SDOH; RWE for breakthrough therapeutics.",
     "3.38"),
    ("3", "Bias-in-translation framework — AoU as pan-NIH diversity benchmark",
     "AoU UBR oversample as a citable diversity reference. ELSI 26.3× + SDOH 5.1× are the political-defense argument. FDA Modernization Act 2.0 creates a regulatory window.",
     "3.75"),
]

t = add_table(s, Inches(0.55), Inches(1.40), Inches(12.25), Inches(4.7),
              header=["#", "Recommendation", "Evidence", "Rubric mean"],
              rows=rows,
              header_font_size=12, body_font_size=11.5,
              col_widths=[0.05, 0.35, 0.50, 0.10],
              last_col_emph=True)
t.rows[0].height = Inches(0.42)
for ri in range(1, 4):
    t.rows[ri].height = Inches(1.42)

add_rect(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.75),
         fill=LIGHT, line=NAVY)
add_textbox(s, Inches(0.75), Inches(6.33), Inches(11.85), Inches(0.6),
            "WG-2 also-considered (deprioritized): NCATS-Translator data linkage; partner-cohort replication pipeline; RWE for breakthrough therapeutics (folded into Rec 2 as use case).",
            size=11, color=DARK, italic=True, font=BODY_FONT, line_spacing=1.25)

add_footer(s, 13)
set_notes(s, """Source: phase4-recommendations/wg2_translational.md.

Rec 1 (3.50): TSBM = Translational Science Benefits Model (Luke et al., Eval Program Plan 2018). Citation/altmetric/RePORTER/NCT APIs all exist; curation and editorial work is new.

Rec 2 (3.38): the named portfolio is the canonical translational-WG deliverable. PGx is the highest peer-index AoU-distinctive use case (4.6×); cardiometabolic+SDOH is the largest combined volume (382 disease pubs + 458 SDOH).

Rec 3 (3.75 — highest WG-2 score): the bias framework is the keystone of the cross-WG hierarchy (WG-2 normative → WG-3 imaging → WG-4 AI). Sets the AoU bias-reporting standard.""")

# ---------------------------------------------------------------------------
# SLIDE 14 — WG-3 Imaging Top-3
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "WG-3 Imaging: starting from infrastructure, not utilization",
           "Lead: Burnside & Giannini. The 2026 Imaging WG begins with 8 corpus pubs — the question is data release, not research demand.")

rows = [
    ("1", "Staged imaging release roadmap by modality",
     "Tier 1: retinal photos, ECG/spirometry waveforms (lowest friction). Tier 2: EHR DICOM (chest X-ray, CT, mammography). Tier 3: prospective MRI in subcohort.",
     "3.25"),
    ("2", "DICOM + FHIR-ImagingStudy standards; radiomics-to-phecode crosswalk",
     "Anchor imaging features to the existing phecode analytic spine (Smoller's 3.80 rubric). NLM-G1 + NIBIB-G4 alignment.",
     "3.88"),
    ("3", "Pre-launch data-quality + bias-audit framework",
     "Inherits the WG-2 bias-in-translation framework. Image-QC per modality; UBR representation audit; per-image 'bias report card'; governance for derivative datasets.",
     "3.50"),
]

t = add_table(s, Inches(0.55), Inches(1.40), Inches(12.25), Inches(4.7),
              header=["#", "Recommendation", "Evidence", "Rubric mean"],
              rows=rows,
              header_font_size=12, body_font_size=11.5,
              col_widths=[0.05, 0.35, 0.50, 0.10],
              last_col_emph=True)
t.rows[0].height = Inches(0.42)
for ri in range(1, 4):
    t.rows[ri].height = Inches(1.42)

add_rect(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.75),
         fill=LIGHT, line=NAVY)
add_textbox(s, Inches(0.75), Inches(6.33), Inches(11.85), Inches(0.6),
            "Reality check: AoU corpus has 8 imaging pubs (HA3, ~0.6%), almost all using EXTERNAL imaging linked to AoU phenotypes. WG-3 starts from infrastructure.",
            size=11, color=AMBER, italic=True, bold=True, font=BODY_FONT, line_spacing=1.25)

add_footer(s, 14)
set_notes(s, """Source: phase4-recommendations/wg3_imaging.md.

Rec 2 (3.88 — highest WG-3 score): DICOM + FHIR + phecode crosswalk is the highest-leverage standards decision because it anchors imaging to existing AoU analytic infrastructure rather than fragmenting the user base.

Cost reality: imaging is the single most expensive recommendation across all 5 WGs. Tier 1 is engineering; Tier 2 is governance; Tier 3 is the actual cost-bearing decision. Defend the staged approach.

Also-considered: AI-radiomics pilot funding (hand to WG-4); UKB/NIH-IDC partnership (flag as cross-cohort validation pathway under Rec 1 Tier 1).""")

# ---------------------------------------------------------------------------
# SLIDE 15 — WG-4 AI Top-3
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "WG-4 AI Strategies: researchers are already using AI — define what AoU provides",
           "Lead: Sastry. HA4 already has 33 AI pubs (30 post-2023); the question is what AoU builds, integrates, or governs.")

rows = [
    ("1", "AoU Workbench AI co-scientist (cohort-builder + code-generator)",
     "NL-to-OMOP/phecode/PRS cohort definition + analytic skeleton. Builds on phecode + CLAMP infrastructure. Reduces time-to-first-result Workbench friction.",
     "3.88"),
    ("2", "AI governance + bias-mitigation framework",
     "Bias auditing for AoU-derived AI; data-use review (training vs inference); Workbench-native bias-audit template; IRB/ELSI integration for participant-facing AI.",
     "3.50"),
    ("3", "Foundation-model evaluation framework (build-vs-buy by end of 2027)",
     "Structured evaluation of existing EHR foundation models (Med-PaLM, EHR-FM family) against AoU benchmarks. Explicit decision-deadline rather than open question.",
     "3.50"),
]

t = add_table(s, Inches(0.55), Inches(1.40), Inches(12.25), Inches(4.7),
              header=["#", "Recommendation", "Evidence", "Rubric mean"],
              rows=rows,
              header_font_size=12, body_font_size=11.5,
              col_widths=[0.05, 0.35, 0.50, 0.10],
              last_col_emph=True)
t.rows[0].height = Inches(0.42)
for ri in range(1, 4):
    t.rows[ri].height = Inches(1.42)

add_rect(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.75),
         fill=SOFT_AMBER, line=AMBER)
add_textbox(s, Inches(0.75), Inches(6.33), Inches(11.85), Inches(0.6),
            "Rec 1 (3.88) is the HIGHEST RUBRIC MEAN across all 5 WGs. Direct researcher impact + standardizes use of validated phenotypes + simplifies analysis.",
            size=11.5, color=AMBER, italic=True, bold=True, font=BODY_FONT, line_spacing=1.25)

add_footer(s, 15)
set_notes(s, """Source: phase4-recommendations/wg4_ai_strategies.md.

Rec 1 (3.88): highest rubric mean in the entire set. The AI co-scientist is the single most-leveraged researcher-facing investment. Defends against 'hallucinations' risk via explicit-provenance design (every cohort/code block links back to source vocabulary entry). 'Build vs buy' answer: build AoU-specific scaffolding on existing foundation model, not the model itself.

Rec 2 (3.50): inherits normative frame from WG-2's bias-in-translation framework. AI-specific failure modes (fairness-through-unawareness, label noise, miscalibration across subgroups, distribution shift) need AI-WG technical expertise.

Rec 3 (3.50): evaluate, don't build. The field is moving so fast that any premature build will be obsolete on delivery. 2027 decision-deadline.""")

# ---------------------------------------------------------------------------
# SLIDE 16 — WG-5 Communications Top-3
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "WG-5 Scientific Communications: turn the corpus evidence into legible products",
           "Lead: Korf. Converts framework + metrics + distinctive evidence into Sci Com / IC / Congressional / participant products.")

rows = [
    ("1", "Annual AoU Impact Report (built on WG-2's TSBM metrics)",
     "Single externally-citable artifact aggregating translational, scientific, equity, and engagement performance. Publish v1 within 12 months.",
     "3.57"),
    ("2", "Plain-language distinctive-evidence narrative",
     "ELSI 26.3× + SDOH 5.1× + PGx 4.6× + methods 3.0×. Anchors the Sci Com ambassador toolkit and IC-PO briefings. \"AoU scaled up everything around genomics that genomics needs to work.\"",
     "3.60"),
    ("3", "Five audience-tuned dissemination products + Sci Com ambassador toolkit",
     "Researcher digest; per-IC briefs (21 already drafted); Congressional 2-page brief; annual participant report; journalist story-pitch sheets.",
     "3.60"),
]

t = add_table(s, Inches(0.55), Inches(1.40), Inches(12.25), Inches(4.7),
              header=["#", "Recommendation", "Evidence", "Rubric mean"],
              rows=rows,
              header_font_size=12, body_font_size=11.5,
              col_widths=[0.05, 0.35, 0.50, 0.10],
              last_col_emph=True)
t.rows[0].height = Inches(0.42)
for ri in range(1, 4):
    t.rows[ri].height = Inches(1.42)

add_rect(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.75),
         fill=LIGHT, line=NAVY)
add_textbox(s, Inches(0.75), Inches(6.33), Inches(11.85), Inches(0.6),
            "The distinctive-evidence narrative (Rec 2) is the political-defense story under which every other WG's recommendations should be read.",
            size=11.5, color=NAVY, italic=True, bold=True, font=BODY_FONT, line_spacing=1.25)

add_footer(s, 16)
set_notes(s, """Source: phase4-recommendations/wg5_communications.md.

Rec 1 (3.57): Impact Report v1 within 12 months. WG-2 defines metrics; WG-5 produces the report. NIH precedent: BD2K, BRAIN Initiative, Cancer Moonshot all publish annual impact products.

Rec 2 (3.60): the four-strength narrative is the MOST EVIDENCE-DRIVEN recommendation in any WG — peer-index numbers are precise from the corpus analysis. Distinctive ≠ growing; the headline rests on these four genuine over-indexes. Dermatology 6.9× over-index is an artifact (rehearsed answer required).

Rec 3 (3.60): the 21 per-IC briefs in aou-ic-alignment/per_ic_briefs/ are mostly drafted. Productize them rather than rebuilding.""")

# ---------------------------------------------------------------------------
# SLIDE 17 — Cross-WG interlocks
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Five cross-WG dependencies the Sci Com should coordinate",
           "Without these handshakes the WGs ship semi-aligned outputs that compete instead of compounding.")

deps = [
    ("CC-1", "9th focus area enables WG-3 + WG-4",
     "WG-1's framework change is the structural prerequisite for both Imaging and AI work to have a coherent home."),
    ("CC-2", "Bias-framework hierarchy across WG-2 / WG-3 / WG-4",
     "WG-2 normative framework → WG-3 imaging audit → WG-4 AI audit. Three WGs should align on a SINGLE AoU bias-reporting standard."),
    ("CC-3", "WG-2 metrics feed WG-5 Impact Report",
     "WG-2 defines TSBM metrics; WG-5 produces the annual Report. Coordinate cadence so Report v1 ships within 12 months."),
    ("CC-4", "WG-5 narrative requires WG-1 framework",
     "The distinctive-evidence narrative depends on framework being settled. Pre-Sci Com circulation should sequence WG-1 → WG-5."),
    ("CC-5", "Aging/ADRD — no WG owns it cleanly",
     "Narrowing peer gap (AoU 6.0% post-2023, ~peer median 0.84×, trend +4.7 vs UKB +6.9 / MVP +5.4 / CKB +5.1). FLAG for WG leads; consider Tier-2 priority + WG-5 patient story."),
]

y0 = Inches(1.42)
row_h = Inches(1.08)
for i, (code, headline, detail) in enumerate(deps):
    y = y0 + row_h * i
    # code box
    add_rect(s, Inches(0.55), y, Inches(1.1), Inches(0.95), fill=NAVY, line=None)
    add_textbox(s, Inches(0.55), y + Inches(0.05), Inches(1.1), Inches(0.4),
                code, size=15, bold=True, color=WHITE, font=HEAD_FONT,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.55), y + Inches(0.45), Inches(1.1), Inches(0.5),
                "interlock", size=8, color=WHITE, font=BODY_FONT,
                align=PP_ALIGN.CENTER, italic=True)
    # headline
    add_textbox(s, Inches(1.85), y + Inches(0.02), Inches(10.95), Inches(0.4),
                headline, size=14, bold=True, color=NAVY, font=HEAD_FONT)
    # detail
    add_textbox(s, Inches(1.85), y + Inches(0.42), Inches(10.95), Inches(0.6),
                detail, size=11.5, color=DARK, font=BODY_FONT, line_spacing=1.25)

add_footer(s, 17)
set_notes(s, """Source: phase4-recommendations/cross_cutting_synthesis.md §1.

The five interlocks should be in the Sci Com's awareness because they determine whether the WGs ship semi-aligned outputs that compete or genuinely-aligned outputs that compound.

CC-5 (Aging/ADRD) is a flag because no Top-3 currently owns it. AoU is at 6.0% post-2023 (≈peer median, index 0.84×) and growing (+4.7), but still trailing UKB's +6.9 surge in the lecanemab/donanemab era. Structural advantage (EHR + UBR diversity); the gap is narrowing rather than flat, so this is now a weaker-evidence flag than in the prior draft. NIA engagement workstream is the natural answer; outside any current Top-3.""")

# ---------------------------------------------------------------------------
# SLIDE 18 — Caveats
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "Caveats — three things the Sci Com should know before acting",
           "Each caveat is a reason to be careful, not a reason to discount the analysis.")

caveats = [
    ("1",
     "Dermatology over-representation is an ARTIFACT",
     "Dermatology is 13.9% of the corpus and 6.9× the peer median — but driven by 1–2 prolific groups, not a strategic priority. Trend score −4.4 (declining share). Do NOT elevate dermatology in the Refresh even though the share index suggests it. WG-5 ambassador toolkit must have the artifact answer rehearsed."),
    ("2",
     "Peer-biobank counts use PubMed name-search fallback",
     "UKB, FinnGen, MVP, and CKB corpora come from PubMed name-search with same-tagger keyword tagging — now full PubMed censuses (UKB 11,868; FinnGen 3,152; MVP 438; CKB 558). Internally consistent across peers but absolute counts under- and over-count differently. USE THE STATIC SHARE INDEX, not trend differentials, for action-relevant comparisons. MVP corpus is small (438); single-MVP signals need a second-peer corroboration."),
    ("3",
     "AoU pre-2023 baseline is tiny (157 pubs)",
     "Any active post-2023 theme will trend positive in AoU's split simply because the corpus grew ~8.1× from pre- to post-2023. Apparent 'growth' on cardiometabolic, mental health, and cancer is largely a small-baseline artifact — read as 'AoU has built a presence here,' not 'AoU is uniquely outpacing the field.' True growth signals come from the peer-comparison panel, not from the AoU-internal pre/post split."),
]

for i, (num, head, body) in enumerate(caveats):
    y = Inches(1.40) + Inches(1.78) * i
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.78), Inches(0.78))
    circ.fill.solid(); circ.fill.fore_color.rgb = AMBER
    circ.line.color.rgb = AMBER
    tf = circ.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.bold = True
    r.font.size = Pt(20); r.font.color.rgb = WHITE; r.font.name = HEAD_FONT
    # head
    add_textbox(s, Inches(1.6), y + Inches(0.02), Inches(11.2), Inches(0.42),
                head, size=15.5, bold=True, color=AMBER, font=HEAD_FONT,
                line_spacing=1.1)
    # body
    add_textbox(s, Inches(1.6), y + Inches(0.45), Inches(11.2), Inches(1.30),
                body, size=11.5, color=DARK, font=BODY_FONT, line_spacing=1.30)

add_footer(s, 18)
set_notes(s, """Source: aou-research-landscape/README.md; peer-biobanks/emerging_themes.md §caveats; phase4-recommendations/README.md.

The dermatology caveat is NON-NEGOTIABLE to include — if Sci Com or external audiences see 'dermatology peer-index 6.9×' without context, they may push to elevate dermatology, which would mislead the Roadmap. Phase 1 README flagged the 1-2-prolific-groups origin; Phase 3b temporal trends show declining share.

The peer-biobank fallback caveat protects against over-reading single-peer trend differentials. The static share index is what we ranked recommendations on.

The pre-2023 baseline caveat is why we don't use AoU-internal pre/post trends as the headline growth metric — we use the peer-comparative gap analysis instead.""")

# ---------------------------------------------------------------------------
# SLIDE 19 — Next steps / asks
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_header(s, "What we ask the Science Committee to decide",
           "Four concrete decisions today; one calendar commitment for the cross-WG hierarchy.")

asks = [
    ("ENDORSE",
     "the addition of a 9th focus area: \"Data, AI & Methods\"",
     "Frame-defining for WG-3 (Imaging) and WG-4 (AI). Without it, downstream WG work is permanently orphaned.",
     ACCENT),
    ("ENDORSE",
     "the FA8 rename to \"Genetics, Multi-omics & Biology\" + commission the proteomics-release roadmap",
     "Signals to the field that the 2026 omics releases have a home. Pair the rename with the data investment.",
     ACCENT),
    ("COMMISSION",
     "the cross-WG bias-in-translation framework hierarchy (WG-2 normative → WG-3 imaging → WG-4 AI)",
     "Single AoU bias-reporting standard rather than three semi-aligned frameworks. WG-2 (Sanchez & Cohn) convenes.",
     NAVY),
    ("ALLOCATE",
     "time at the next Science Committee for the structural-framework vote (FA9 + FA8 rename + FA reframes)",
     "These three structural decisions are the WG-1 deliverable; they should be settled before the other four WGs' Top-3 are finalized.",
     NAVY),
]

y0 = Inches(1.40)
ask_h = Inches(1.20)
for i, (action, headline, detail, color) in enumerate(asks):
    y = y0 + ask_h * i
    # left action chip
    add_rect(s, Inches(0.55), y, Inches(2.1), Inches(1.05), fill=color, line=None)
    add_textbox(s, Inches(0.55), y + Inches(0.25), Inches(2.1), Inches(0.55),
                action, size=18, bold=True, color=WHITE, font=HEAD_FONT,
                align=PP_ALIGN.CENTER)
    # headline + detail
    add_textbox(s, Inches(2.85), y + Inches(0.05), Inches(9.95), Inches(0.5),
                headline, size=14, bold=True, color=NAVY, font=HEAD_FONT,
                line_spacing=1.12)
    add_textbox(s, Inches(2.85), y + Inches(0.55), Inches(9.95), Inches(0.55),
                detail, size=11.5, color=DARK, font=BODY_FONT, italic=True,
                line_spacing=1.25)

# Closing line
add_rect(s, Inches(0.55), Inches(6.30), Inches(12.25), Inches(0.65),
         fill=NAVY, line=None)
add_textbox(s, Inches(0.75), Inches(6.36), Inches(11.85), Inches(0.55),
            "Thank you. We look forward to the Committee's discussion.   —  Jordan Smoller & Melissa Haendel, WG-1 Co-leads",
            size=12.5, bold=True, color=WHITE, font=HEAD_FONT, italic=True,
            align=PP_ALIGN.CENTER)

add_footer(s, 19)
set_notes(s, """The four explicit asks the Sci Com should take action on:

(1) ENDORSE the 9th focus area — frame-defining; can only be done by this WG.
(2) ENDORSE the FA8 rename + commission the proteomics roadmap — pair the structural signal with the data investment.
(3) COMMISSION the cross-WG bias framework hierarchy — WG-2 (Sanchez & Cohn) convenes WG-3 + WG-4.
(4) ALLOCATE time at the next Sci Com for the structural-framework vote — sequence WG-1 decisions BEFORE the other four WGs' Top-3 are finalized.

Smaller asks tucked in: surface Aging/ADRD as a Tier-2 priority for WG-5 patient-facing story.

If Sci Com only takes one action: endorse the 9th focus area. It's the structural prerequisite for everything else.""")

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Slides: {len(prs.slides)}")
